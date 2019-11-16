import zipfile
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import time
import os
import json
import traceback

start_time = time.time()

script_dir = os.path.dirname(__file__)
configuration = json.load(open(os.path.join(script_dir, 'config.json'),"r"))
project_locations = configuration["files_for_processing"]


number_of_doc_processed = 0
number_of_ppt_processed = 0
number_of_exceptions = 0
number_of_docx_img_extracted = 0
number_of_pptx_img_extracted = 0
log = ""


def iter_picture_shapes(prs):
    for index, slide in enumerate(prs.slides):
        img_no = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_no = img_no + 1
                yield [shape, index+1, img_no]

def pptx_img_extractor(root, file):
    global number_of_pptx_img_extracted, number_of_exceptions
    path =  os.path.join(root, file)
    try:
        for index, picture in enumerate(iter_picture_shapes(Presentation(path))):
            image = picture[0].image
            # ---get image "file" contents---
            image_bytes = image.blob
            image_filename = 'img_' + str(picture[1]) + '_' + str(picture[2]) + '.' + image.ext
            with open(os.path.join(root, image_filename), 'wb') as f:
                f.write(image_bytes)
                number_of_pptx_img_extracted = number_of_pptx_img_extracted + 1
    except Exception as e:
        number_of_exceptions = number_of_exceptions + 1
        print("Exception while processing pptx file - " + path)
        traceback.print_tb(e.__traceback__)
        print(traceback.print_tb(e.__traceback__))

def doc_img_extractor_zip(root, file):
    global number_of_docx_img_extracted, number_of_exceptions
    path = os.path.join(root, file)
    try:
        z = zipfile.ZipFile(path)

        all_files = z.namelist()
        # get all files in word/media/ directory
        images = list(filter(lambda x: x.startswith('word/media/'), all_files))

        for index, image in enumerate(images):
            # open an image and save it
            img = z.open(image).read()

            img_name = "img_" + str(index + 1) + "." + get_file_extension(image)
            f = open(os.path.join(root, img_name), 'wb')
            f.write(img)
            number_of_docx_img_extracted = number_of_docx_img_extracted + 1
    except Exception as e:
        number_of_exceptions = number_of_exceptions + 1
        print("Exception while processing docx file - " + path)
        traceback.print_tb(e.__traceback__)
        print(traceback.print_tb(e.__traceback__))


# def doc_img_extractor_docx(root, file):
#     global number_of_docx_img_extracted, number_of_exceptions
#     path = os.path.join(root, file)
#     try:
#         doc = docx.Document(path)
#         for index, s in enumerate(doc.inline_shapes):
#             name = s._inline.graphic.graphicData.pic.nvPicPr.cNvPr.name
#             img_name = "img_" + str(index + 1) + get_file_extension(img)
#             f = open(os.path.join(root, img_name), 'wb')
#             f.write(img)
#             number_of_docx_img_extracted = number_of_docx_img_extracted + 1
#     except Exception as e:
#         number_of_exceptions = number_of_exceptions + 1
#         print("Exception while processing docx file - " + path)
#         traceback.print_tb(e.__traceback__)
#         log.error(traceback.print_tb(e.__traceback__))


def get_file_extension(file_name):
    if len(file_name.split(os.extsep,1)) < 2:
        ext = 'no format'
    else:
        ext = file_name.split(os.extsep,1)[-1]
    return ext


def main():
    global number_of_doc_processed, number_of_ppt_processed, number_of_exceptions

    for project_location in project_locations:
        print(project_location["path"])
        try:
            for root, dirs, files in os.walk(project_location["path"], topdown=True):
                for file in files:
                    file_ext = get_file_extension(file)
                    if file_ext in ".docx":
                        doc_img_extractor_zip(root, file)
                        number_of_doc_processed = number_of_doc_processed + 1
                    elif file_ext in ".pptx":
                        pptx_img_extractor(root, file)
                        number_of_ppt_processed = number_of_ppt_processed + 1
        except Exception as e:
            number_of_exceptions = number_of_exceptions + 1
            print("Exception while processing dir - " + project_location["path"])
            traceback.print_tb(e.__traceback__)
            print(traceback.print_tb(e.__traceback__))


main()
print("Total number of doc file processed: " + str (number_of_doc_processed))
print("Total number of ppt file processed: " + str (number_of_ppt_processed))
print("Number of images extracted from docx files: " + str (number_of_docx_img_extracted))
print("Number of images extracted from pptx files: " + str (number_of_pptx_img_extracted))
print("Number of exceptions occurred while processing - " + str(number_of_exceptions))
print("Total time took for the execution - "+ str((time.time() - start_time)/60) + " minutes")